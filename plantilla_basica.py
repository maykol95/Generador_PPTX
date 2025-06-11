from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from io import BytesIO
import os
import time

def calcular_dimensiones(img_path, fotos_por_slide, encabezados_count):
    try:
        with Image.open(img_path) as img:
            width_px, height_px = img.size
            width_in, height_in = width_px / 96, height_px / 96
            max_width = 11 / fotos_por_slide - 0.5
            encabezado_space = 0.3 + 0.25 * min(encabezados_count, 6)
            max_height = 5.3 - encabezado_space
            ratio = min(max_width / width_in, max_height / height_in)
            return Inches(width_in * ratio), Inches(height_in * ratio)
    except Exception as e:
        print(f"Error al procesar {img_path}: {e}")
        return None, None

def generar_presentacion_basica(
    df_slice,
    nombre,
    fotos_por_slide,
    encabezados_seleccionados,
    fondo_bytes,
    tipo_fuente,
    color_fuente,
    temp_dir,
    tiempo_texto=None,
    start_time=None
):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    fondo_path = None
    if fondo_bytes:
        try:
            fondo_img = Image.open(BytesIO(fondo_bytes.read()))
            fondo_path = os.path.join(temp_dir, "fondo_temp.jpg")
            fondo_img.save(fondo_path)
        except Exception as e:
            print(f"Error al cargar fondo: {e}")
            fondo_path = None

    for i in range(0, len(df_slice), fotos_por_slide):
        slide = prs.slides.add_slide(layout)
        if fondo_path:
            try:
                slide.shapes.add_picture(fondo_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            except Exception as e:
                print(f"Error al insertar fondo: {e}")

        subset = df_slice.iloc[i:i+fotos_por_slide]
        spacing_x = Inches(12 / fotos_por_slide)

        for j, (_, row) in enumerate(subset.iterrows()):
            img_path = row['img_path']
            img_width, img_height = calcular_dimensiones(img_path, fotos_por_slide, len(encabezados_seleccionados))
            img_left = spacing_x * j + Inches(1.15)
            encabezado_height = Inches(0.25 * min(len(encabezados_seleccionados), 6))
            img_top = encabezado_height + Inches(0.4)

            if img_width and img_height:
                try:
                    slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
                except Exception as e:
                    print(f"Error al insertar imagen {img_path}: {e}")

                encabezado = " | ".join(f"{col}: {row.get(col, '')}" for col in encabezados_seleccionados)
                if encabezado:
                    try:
                        text_box = slide.shapes.add_textbox(img_left, Inches(0.2), img_width, encabezado_height)
                        text_frame = text_box.text_frame
                        text_frame.clear()
                        for line in encabezado.split(" | "):
                            p = text_frame.add_paragraph()
                            p.text = line
                            p.font.size = Pt(8)
                            p.font.bold = True
                            p.font.name = tipo_fuente
                            r, g, b = [int(color_fuente.lstrip("#")[k:k+2], 16) for k in (0, 2, 4)]
                            p.font.color.rgb = RGBColor(r, g, b)
                    except Exception as e:
                        print(f"Error al agregar encabezado: {e}")

        if tiempo_texto and start_time:
            tiempo_actual = round(time.time() - start_time, 2)
            tiempo_texto.markdown(f"⏳ Tiempo transcurrido: **{tiempo_actual} segundos**")

    path = os.path.join(temp_dir, f"{nombre}.pptx")
    prs.save(path)
    return [(nombre, path)]
