from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
from io import BytesIO
import os
import time

def generar_presentacion_exhibiciones(
    df,
    agrupador_1,
    agrupador_2,
    fondo_bytes,
    tipo_fuente,
    status_text,
    progress_bar,
    temp_dir,
    columna_img,
    color_agrupador_1,
    color_agrupador_2,
    color_encabezado,
    columna_encabezado=None,
    tiempo_texto=None,
    start_time=None
):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    fondo_path = None
    if fondo_bytes:
        try:
            fondo_img = Image.open(BytesIO(fondo_bytes.read()))
            fondo_path = os.path.join(temp_dir, "fondo_temp.jpg")
            fondo_img.save(fondo_path)
        except Exception as e:
            print(f"Error al cargar fondo: {e}")
            fondo_path = None

    grupos_principales = df[agrupador_1].dropna().unique()

    r1, g1, b1 = [int(color_agrupador_1.lstrip("#")[i:i+2], 16) for i in (0, 2, 4)]
    r2, g2, b2 = [int(color_agrupador_2.lstrip("#")[i:i+2], 16) for i in (0, 2, 4)]
    r3, g3, b3 = [int(color_encabezado.lstrip("#")[i:i+2], 16) for i in (0, 2, 4)]

    img_width, img_height = Inches(2), Inches(2)
    margin_x, margin_y = Inches(2), Inches(0.5)
    spacing_x, spacing_y = Inches(2.2), Inches(0.5)
    max_x = prs.slide_width - margin_x
    max_y = prs.slide_height - margin_y

    def nueva_slide():
        slide = prs.slides.add_slide(layout)
        if fondo_path:
            try:
                slide.shapes.add_picture(fondo_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            except Exception as e:
                print(f"Error al insertar fondo: {e}")
        return slide

    def agregar_titulo(slide, titulo):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.5), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = str(titulo)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = tipo_fuente
        p.font.color.rgb = RGBColor(r1, g1, b1)

    def agregar_subtitulo(slide, subtitulo, y):
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.5), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = str(subtitulo)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.name = tipo_fuente
        p.font.color.rgb = RGBColor(r2, g2, b2)

    def agregar_encabezado(slide, texto, x, y):
        box = slide.shapes.add_textbox(x, y, img_width, Inches(0.25))
        p = box.text_frame.paragraphs[0]
        p.text = str(texto)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.name = tipo_fuente
        p.font.color.rgb = RGBColor(r3, g3, b3)

    for idx_grupo1, grupo_1 in enumerate(grupos_principales):
        df_grupo = df[df[agrupador_1] == grupo_1]
        subgrupos = df_grupo[agrupador_2].dropna().unique()

        current_slide = nueva_slide()
        agregar_titulo(current_slide, grupo_1)
        current_y = Inches(0.8)

        for subgrupo in subgrupos:
            df_subgrupo = df_grupo[df_grupo[agrupador_2] == subgrupo]

            required_height = Inches(0.4) + img_height
            if current_y + required_height > max_y:
                current_slide = nueva_slide()
                agregar_titulo(current_slide, grupo_1)
                current_y = Inches(0.8)

            agregar_subtitulo(current_slide, subgrupo, current_y)
            current_y += Inches(0.4)

            row_x = margin_x
            row_y = current_y

            for _, row in df_subgrupo.iterrows():
                img_path = row.get(columna_img)
                encabezado = row.get(columna_encabezado) if columna_encabezado else None

                if not img_path or not os.path.exists(img_path):
                    print(f"Imagen no encontrada: {img_path}")
                    continue

                if row_x + img_width > max_x:
                    row_x = margin_x
                    row_y += img_height + spacing_y

                extra_y = Inches(0.3) if encabezado else 0

                if row_y + img_height + extra_y > max_y:
                    current_slide = nueva_slide()
                    agregar_titulo(current_slide, grupo_1)
                    agregar_subtitulo(current_slide, subgrupo, Inches(0.8))
                    row_x = margin_x
                    row_y = Inches(1.2)

                if encabezado:
                    agregar_encabezado(current_slide, encabezado, row_x, row_y)
                    row_y += Inches(0.25)

                try:
                    current_slide.shapes.add_picture(img_path, row_x, row_y, width=img_width, height=img_height)
                except Exception as e:
                    print(f"Error al insertar imagen: {e}")

                row_x += spacing_x
                row_y = row_y - Inches(0.25) if encabezado else row_y

            current_y = row_y + img_height + spacing_y

        progress_bar.progress((idx_grupo1 + 1) / len(grupos_principales))
        status_text.text(f"Grupo {idx_grupo1 + 1} de {len(grupos_principales)} procesado")

        if tiempo_texto and start_time:
            tiempo_actual = round(time.time() - start_time, 2)
            tiempo_texto.markdown(f"⏳ Tiempo transcurrido: **{tiempo_actual} segundos**")

    path = os.path.join(temp_dir, "presentacion_exhibiciones_completa.pptx")
    prs.save(path)
    return [("presentacion_exhibiciones_completa", path)]
